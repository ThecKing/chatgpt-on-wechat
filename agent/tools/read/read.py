# -*- coding: utf-8 -*-
"""
读取文件工具

支持读取文本文件、图片、PDF等多种格式：
- 文本文件：返回内容（自动截断）
- 图片文件：返回元数据（用于发送）
- PDF文件：提取文本内容
- Office文件：提取文本内容
- 音视频/二进制文件：返回元数据
"""

# 导入操作系统模块
import os

# 导入类型提示
from typing import Dict, Any

# 导入路径处理
from pathlib import Path

# 导入工具基类和结果类
from agent.tools.base_tool import BaseTool, ToolResult

# 导入截断工具
from agent.tools.utils.truncate import truncate_head, format_size, DEFAULT_MAX_LINES, DEFAULT_MAX_BYTES

# 导入路径扩展工具
from common.utils import expand_path


class Read(BaseTool):
    """
    读取文件工具类
    
    支持多种文件类型的读取：
    - 文本文件：返回内容（带截断）
    - 图片：返回元数据
    - PDF：提取文本
    - Office：提取文本
    - 其他：返回元数据
    """
    
    # 工具名称
    name: str = "read"
    
    # 工具描述
    description: str = f"Read or inspect file contents. For text/PDF files, returns content (truncated to {DEFAULT_MAX_LINES} lines or {DEFAULT_MAX_BYTES // 1024}KB). For images/videos/audio, returns metadata only (file info, size, type). Use offset/limit for large text files."
    
    # 参数JSON Schema
    params: dict = {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the file to read. IMPORTANT: Relative paths are based on workspace directory. To access files outside workspace, use absolute paths starting with ~ or /."
            },
            "offset": {
                "type": "integer",
                "description": "Line number to start reading from (1-indexed, optional). Use negative values to read from end (e.g. -20 for last 20 lines)"
            },
            "limit": {
                "type": "integer",
                "description": "Maximum number of lines to read (optional)"
            }
        },
        "required": ["path"]
    }
    
    def __init__(self, config: dict = None):
        """
        初始化读取工具
        
        Args:
            config: 配置字典
        """
        # 存储配置
        self.config = config or {}
        # 获取工作目录
        self.cwd = self.config.get("cwd", os.getcwd())
        
        # 定义图片扩展名集合
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg', '.ico'}
        # 定义视频扩展名集合
        self.video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.webm', '.m4v'}
        # 定义音频扩展名集合
        self.audio_extensions = {'.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma'}
        # 定义二进制扩展名集合
        self.binary_extensions = {'.exe', '.dll', '.so', '.dylib', '.bin', '.dat', '.db', '.sqlite'}
        # 定义压缩包扩展名集合
        self.archive_extensions = {'.zip', '.tar', '.gz', '.rar', '.7z', '.bz2', '.xz'}
        # 定义PDF扩展名集合
        self.pdf_extensions = {'.pdf'}
        # 定义Office扩展名集合
        self.office_extensions = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}

        # 定义可读文本格式扩展名集合
        self.text_extensions = {
            # 纯文本和标记语言
            '.txt', '.md', '.markdown', '.rst', '.log', '.csv', '.tsv', '.json', '.xml', '.yaml', '.yml',
            # 编程语言源代码
            '.py', '.js', '.ts', '.java', '.c', '.cpp', '.h', '.hpp', '.go', '.rs', '.rb', '.php',
            # Web前端
            '.html', '.css', '.scss', '.sass', '.less', '.vue', '.jsx', '.tsx',
            # 脚本语言
            '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
            # 其他编程语言和配置
            '.sql', '.r', '.m', '.swift', '.kt', '.scala', '.clj', '.erl', '.ex',
            # 构建和配置文件
            '.dockerfile', '.makefile', '.cmake', '.gradle', '.properties', '.ini', '.conf', '.cfg',
        }
    
    def execute(self, args: Dict[str, Any]) -> ToolResult:
        """
        执行文件读取操作
        
        Args:
            args: 包含文件路径和可选的offset/limit参数
            
        Returns:
            ToolResult: 文件内容或错误信息
        """
        # 支持 'location' 作为 'path' 的别名（LLM可能从skill列表中使用）
        path = args.get("path", "") or args.get("location", "")
        path = path.strip() if isinstance(path, str) else ""
        offset = args.get("offset")       # 起始行号
        limit = args.get("limit")         # 最大行数

        # 检查路径参数
        if not path:
            return ToolResult.fail("Error: path parameter is required")
        
        # 解析路径
        absolute_path = self._resolve_path(path)
        
        # 安全检查：防止读取敏感配置文件
        env_config_path = expand_path("~/.cow/.env")
        if os.path.abspath(absolute_path) == os.path.abspath(env_config_path):
            return ToolResult.fail(
                "Error: Access denied. API keys and credentials must be accessed through the env_config tool only."
            )
        
        # 检查文件是否存在
        if not os.path.exists(absolute_path):
            # 如果使用相对路径，提供有用提示
            if not os.path.isabs(path) and not path.startswith('~'):
                return ToolResult.fail(
                    f"Error: File not found: {path}\n"
                    f"Resolved to: {absolute_path}\n"
                    f"Hint: Relative paths are based on workspace ({self.cwd}). For files outside workspace, use absolute paths."
                )
            return ToolResult.fail(f"Error: File not found: {path}")
        
        # 检查是否可读
        if not os.access(absolute_path, os.R_OK):
            return ToolResult.fail(f"Error: File is not readable: {path}")
        
        # 获取文件扩展名和大小
        file_ext = Path(absolute_path).suffix.lower()
        file_size = os.path.getsize(absolute_path)
        
        # 如果是图片 - 返回元数据用于发送
        if file_ext in self.image_extensions:
            return self._read_image(absolute_path, file_ext)
        
        # 如果是视频/音频/二进制/压缩包 - 只返回元数据
        if file_ext in self.video_extensions:
            return self._return_file_metadata(absolute_path, "video", file_size)
        if file_ext in self.audio_extensions:
            return self._return_file_metadata(absolute_path, "audio", file_size)
        if file_ext in self.binary_extensions or file_ext in self.archive_extensions:
            return self._return_file_metadata(absolute_path, "binary", file_size)
        
        # 如果是PDF
        if file_ext in self.pdf_extensions:
            return self._read_pdf(absolute_path, path, offset, limit)

        # 如果是Office文档
        if file_ext in self.office_extensions:
            return self._read_office(absolute_path, path, file_ext, offset, limit)

        # 读取文本文件（大文件自动截断）
        return self._read_text(absolute_path, path, offset, limit)
    
    def _resolve_path(self, path: str) -> str:
        """
        解析路径为绝对路径
        
        Args:
            path: 相对或绝对路径
            
        Returns:
            str: 绝对路径
        """
        # 扩展 ~ 为用户主目录
        path = expand_path(path)
        if os.path.isabs(path):
            return path
        # 相对路径基于工作目录
        return os.path.abspath(os.path.join(self.cwd, path))
    
    def _return_file_metadata(self, absolute_path: str, file_type: str, file_size: int) -> ToolResult:
        """
        返回不可读文件的元数据（视频、音频、二进制等）
        
        Args:
            absolute_path: 文件的绝对路径
            file_type: 文件类型
            file_size: 文件大小（字节）
            
        Returns:
            ToolResult: 文件元数据
        """
        # 获取文件名和扩展名
        file_name = Path(absolute_path).name
        file_ext = Path(absolute_path).suffix.lower()
        
        # MIME类型映射
        mime_types = {
            # 视频
            '.mp4': 'video/mp4', '.avi': 'video/x-msvideo', '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska', '.webm': 'video/webm',
            # 音频
            '.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.ogg': 'audio/ogg',
            '.m4a': 'audio/mp4', '.flac': 'audio/flac',
            # 二进制
            '.zip': 'application/zip', '.tar': 'application/x-tar',
            '.gz': 'application/gzip', '.rar': 'application/x-rar-compressed',
        }
        # 获取MIME类型，默认为二进制流
        mime_type = mime_types.get(file_ext, 'application/octet-stream')
        
        # 构建结果
        result = {
            "type": f"{file_type}_metadata",   # 结果类型
            "file_type": file_type,             # 文件类型
            "path": absolute_path,              # 绝对路径
            "file_name": file_name,             # 文件名
            "mime_type": mime_type,             # MIME类型
            "size": file_size,                  # 文件大小
            "size_formatted": format_size(file_size),  # 格式化大小
            "message": f"{file_type.capitalize()} 文件: {file_name} ({format_size(file_size)})\n提示: 如果需要发送此文件，请使用 send 工具。"
        }
        
        return ToolResult.success(result)
    
    def _read_image(self, absolute_path: str, file_ext: str) -> ToolResult:
        """
        读取图片文件 - 只返回元数据（图片应该发送而非读取到上下文）
        
        Args:
            absolute_path: 图片文件的绝对路径
            file_ext: 文件扩展名
            
        Returns:
            ToolResult: 包含图片元数据的结果
        """
        try:
            # 获取文件大小
            file_size = os.path.getsize(absolute_path)
            
            # MIME类型映射
            mime_type_map = {
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.png': 'image/png',
                '.gif': 'image/gif',
                '.webp': 'image/webp'
            }
            mime_type = mime_type_map.get(file_ext, 'image/jpeg')
            
            # 返回元数据（不是file_to_send - 使用send工具实际发送）
            result = {
                "type": "image_metadata",       # 结果类型
                "file_type": "image",           # 文件类型
                "path": absolute_path,          # 绝对路径
                "mime_type": mime_type,         # MIME类型
                "size": file_size,              # 文件大小
                "size_formatted": format_size(file_size),  # 格式化大小
                "message": f"图片文件: {Path(absolute_path).name} ({format_size(file_size)})\n提示: 如果需要发送此图片，请使用 send 工具。"
            }
            
            return ToolResult.success(result)
            
        except Exception as e:
            return ToolResult.fail(f"Error reading image file: {str(e)}")
    
    def _read_text(self, absolute_path: str, display_path: str, offset: int = None, limit: int = None) -> ToolResult:
        """
        读取文本文件
        
        Args:
            absolute_path: 文件的绝对路径
            display_path: 显示用的路径
            offset: 起始行号（1索引）
            limit: 最大读取行数
            
        Returns:
            ToolResult: 文件内容或错误信息
        """
        try:
            # 首先检查文件大小
            file_size = os.path.getsize(absolute_path)
            MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
            
            # 如果文件过大，返回元数据
            if file_size > MAX_FILE_SIZE:
                return ToolResult.success({
                    "type": "file_to_send",
                    "file_type": "document",
                    "path": absolute_path,
                    "size": file_size,
                    "size_formatted": format_size(file_size),
                    "message": f"文件过大 ({format_size(file_size)} > 50MB)，无法读取内容。文件路径: {absolute_path}"
                })
            
            # 读取文件（utf-8-sig自动去除Windows BOM）
            with open(absolute_path, 'r', encoding='utf-8-sig') as f:
                content = f.read()
            
            # 如果内容过长则截断（模型上下文最多20K字符）
            MAX_CONTENT_CHARS = 20 * 1024  # 20K字符
            content_truncated = False
            if len(content) > MAX_CONTENT_CHARS:
                content = content[:MAX_CONTENT_CHARS]
                content_truncated = True
            
            # 分割为行
            all_lines = content.split('\n')
            total_file_lines = len(all_lines)
            
            # 应用offset（如果指定）
            start_line = 0
            if offset is not None:
                if offset < 0:
                    # 负偏移：从末尾读取
                    # -20 表示"最后20行" -> 从 (total - 20) 开始
                    start_line = max(0, total_file_lines + offset)
                else:
                    # 正偏移：从开头读取（1索引）
                    start_line = max(0, offset - 1)  # 转换为0索引
                    if start_line >= total_file_lines:
                        return ToolResult.fail(
                            f"Error: Offset {offset} is beyond end of file ({total_file_lines} lines total)"
                        )
            
            # 显示用的起始行号（1索引）
            start_line_display = start_line + 1
            
            # 如果用户指定了limit，使用它
            selected_content = content
            user_limited_lines = None
            if limit is not None:
                end_line = min(start_line + limit, total_file_lines)
                selected_content = '\n'.join(all_lines[start_line:end_line])
                user_limited_lines = end_line - start_line
            elif offset is not None:
                selected_content = '\n'.join(all_lines[start_line:])
            
            # 应用截断（考虑行数和字节限制）
            truncation = truncate_head(selected_content)
            
            # 初始化输出文本和详情
            output_text = ""
            details = {}
            
            # 如果内容被截断，添加警告
            if content_truncated:
                output_text = f"[文件内容已截断到前 {format_size(MAX_CONTENT_CHARS)}，完整文件大小: {format_size(file_size)}]\n\n"
            
            if truncation.first_line_exceeds_limit:
                # 第一行超过30KB限制
                first_line_size = format_size(len(all_lines[start_line].encode('utf-8')))
                output_text = f"[Line {start_line_display} is {first_line_size}, exceeds {format_size(DEFAULT_MAX_BYTES)} limit. Use bash tool to read: head -c {DEFAULT_MAX_BYTES} {display_path} | tail -n +{start_line_display}]"
                details["truncation"] = truncation.to_dict()
            elif truncation.truncated:
                # 发生了截断
                end_line_display = start_line_display + truncation.output_lines - 1
                next_offset = end_line_display + 1
                
                output_text = truncation.content
                
                if truncation.truncated_by == "lines":
                    output_text += f"\n\n[Showing lines {start_line_display}-{end_line_display} of {total_file_lines}. Use offset={next_offset} to continue.]"
                else:
                    output_text += f"\n\n[Showing lines {start_line_display}-{end_line_display} of {total_file_lines} ({format_size(DEFAULT_MAX_BYTES)} limit). Use offset={next_offset} to continue.]"
                
                details["truncation"] = truncation.to_dict()
            elif user_limited_lines is not None and start_line + user_limited_lines < total_file_lines:
                # 用户指定了limit，还有更多内容可用
                remaining = total_file_lines - (start_line + user_limited_lines)
                next_offset = start_line + user_limited_lines + 1
                
                output_text = truncation.content
                output_text += f"\n\n[{remaining} more lines in file. Use offset={next_offset} to continue.]"
            else:
                # 没有截断，没有超过用户限制
                output_text = truncation.content
            
            # 构建结果
            result = {
                "content": output_text,           # 文件内容
                "total_lines": total_file_lines,  # 总行数
                "start_line": start_line_display, # 起始行号
                "output_lines": truncation.output_lines  # 输出行数
            }
            
            # 如果有详情，添加到结果
            if details:
                result["details"] = details
            
            return ToolResult.success(result)
            
        except UnicodeDecodeError:
            return ToolResult.fail(f"Error: File is not a valid text file (encoding error): {display_path}")
        except Exception as e:
            return ToolResult.fail(f"Error reading file: {str(e)}")
    
    def _read_office(self, absolute_path: str, display_path: str, file_ext: str,
                     offset: int = None, limit: int = None) -> ToolResult:
        """
        读取Office文档（.docx, .xlsx, .pptx）
        
        使用 python-docx / openpyxl / python-pptx 库提取文本。
        
        Args:
            absolute_path: 文件绝对路径
            display_path: 显示路径
            file_ext: 文件扩展名
            offset: 起始行号
            limit: 最大行数
            
        Returns:
            ToolResult: 提取的文本内容
        """
        try:
            # 提取Office文档文本
            text = self._extract_office_text(absolute_path, file_ext)
        except ImportError as e:
            return ToolResult.fail(str(e))
        except Exception as e:
            return ToolResult.fail(f"Error reading Office document: {e}")

        # 如果没有提取到文本
        if not text or not text.strip():
            return ToolResult.success({
                "content": f"[Office file {Path(absolute_path).name}: no text content could be extracted]",
            })

        # 分割为行
        all_lines = text.split('\n')
        total_lines = len(all_lines)

        # 计算起始行
        start_line = 0
        if offset is not None:
            if offset < 0:
                start_line = max(0, total_lines + offset)
            else:
                start_line = max(0, offset - 1)
                if start_line >= total_lines:
                    return ToolResult.fail(
                        f"Error: Offset {offset} is beyond end of content ({total_lines} lines total)"
                    )

        # 选择内容
        selected_content = text
        user_limited_lines = None
        if limit is not None:
            end_line = min(start_line + limit, total_lines)
            selected_content = '\n'.join(all_lines[start_line:end_line])
            user_limited_lines = end_line - start_line
        elif offset is not None:
            selected_content = '\n'.join(all_lines[start_line:])

        # 应用截断
        truncation = truncate_head(selected_content)
        start_line_display = start_line + 1
        output_text = ""

        if truncation.truncated:
            end_line_display = start_line_display + truncation.output_lines - 1
            next_offset = end_line_display + 1
            output_text = truncation.content
            output_text += f"\n\n[Showing lines {start_line_display}-{end_line_display} of {total_lines}. Use offset={next_offset} to continue.]"
        elif user_limited_lines is not None and start_line + user_limited_lines < total_lines:
            remaining = total_lines - (start_line + user_limited_lines)
            next_offset = start_line + user_limited_lines + 1
            output_text = truncation.content
            output_text += f"\n\n[{remaining} more lines in file. Use offset={next_offset} to continue.]"
        else:
            output_text = truncation.content

        return ToolResult.success({
            "content": output_text,
            "total_lines": total_lines,
            "start_line": start_line_display,
            "output_lines": truncation.output_lines,
        })

    @staticmethod
    def _extract_office_text(absolute_path: str, file_ext: str) -> str:
        """
        从Office文档提取纯文本
        
        Args:
            absolute_path: 文件绝对路径
            file_ext: 文件扩展名
            
        Returns:
            str: 提取的文本
            
        Raises:
            ImportError: 如果缺少必要的库
        """
        # 处理Word文档
        if file_ext in ('.docx', '.doc'):
            try:
                from docx import Document
            except ImportError:
                raise ImportError("Error: python-docx library not installed. Install with: pip install python-docx")
            doc = Document(absolute_path)
            paragraphs = [p.text for p in doc.paragraphs]
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    paragraphs.append('\t'.join(cell.text for cell in row.cells))
            return '\n'.join(paragraphs)

        # 处理Excel文档
        if file_ext in ('.xlsx', '.xls'):
            try:
                from openpyxl import load_workbook
            except ImportError:
                raise ImportError("Error: openpyxl library not installed. Install with: pip install openpyxl")
            wb = load_workbook(absolute_path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"--- Sheet: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    parts.append('\t'.join(str(c) if c is not None else '' for c in row))
            wb.close()
            return '\n'.join(parts)

        # 处理PowerPoint文档
        if file_ext in ('.pptx', '.ppt'):
            try:
                from pptx import Presentation
            except ImportError:
                raise ImportError("Error: python-pptx library not installed. Install with: pip install python-pptx")
            prs = Presentation(absolute_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                parts.append(text)
            return '\n'.join(parts)

        return ""

    def _read_pdf(self, absolute_path: str, display_path: str, offset: int = None, limit: int = None) -> ToolResult:
        """
        读取PDF文件内容
        
        Args:
            absolute_path: 文件的绝对路径
            display_path: 显示路径
            offset: 起始行号（1索引）
            limit: 最大读取行数
            
        Returns:
            ToolResult: PDF文本内容或错误信息
        """
        try:
            # 尝试导入pypdf
            try:
                from pypdf import PdfReader
            except ImportError:
                return ToolResult.fail(
                    "Error: pypdf library not installed. Install with: pip install pypdf"
                )
            
            # 创建PDF阅读器
            reader = PdfReader(absolute_path)
            total_pages = len(reader.pages)
            
            # 从所有页面提取文本
            text_parts = []
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")
            
            # 如果没有提取到文本
            if not text_parts:
                return ToolResult.success({
                    "content": f"[PDF file with {total_pages} pages, but no text content could be extracted]",
                    "total_pages": total_pages,
                    "message": "PDF may contain only images or be encrypted"
                })
            
            # 合并所有文本
            full_content = "\n\n".join(text_parts)
            all_lines = full_content.split('\n')
            total_lines = len(all_lines)
            
            # 应用offset和limit（与文本文件相同逻辑）
            start_line = 0
            if offset is not None:
                start_line = max(0, offset - 1)
                if start_line >= total_lines:
                    return ToolResult.fail(
                        f"Error: Offset {offset} is beyond end of content ({total_lines} lines total)"
                    )
            
            start_line_display = start_line + 1
            
            # 选择内容
            selected_content = full_content
            user_limited_lines = None
            if limit is not None:
                end_line = min(start_line + limit, total_lines)
                selected_content = '\n'.join(all_lines[start_line:end_line])
                user_limited_lines = end_line - start_line
            elif offset is not None:
                selected_content = '\n'.join(all_lines[start_line:])
            
            # 应用截断
            truncation = truncate_head(selected_content)
            
            output_text = ""
            details = {}
            
            if truncation.truncated:
                end_line_display = start_line_display + truncation.output_lines - 1
                next_offset = end_line_display + 1
                
                output_text = truncation.content
                
                if truncation.truncated_by == "lines":
                    output_text += f"\n\n[Showing lines {start_line_display}-{end_line_display} of {total_lines}. Use offset={next_offset} to continue.]"
                else:
                    output_text += f"\n\n[Showing lines {start_line_display}-{end_line_display} of {total_lines} ({format_size(DEFAULT_MAX_BYTES)} limit). Use offset={next_offset} to continue.]"
                
                details["truncation"] = truncation.to_dict()
            elif user_limited_lines is not None and start_line + user_limited_lines < total_lines:
                remaining = total_lines - (start_line + user_limited_lines)
                next_offset = start_line + user_limited_lines + 1
                
                output_text = truncation.content
                output_text += f"\n\n[{remaining} more lines in file. Use offset={next_offset} to continue.]"
            else:
                output_text = truncation.content
            
            # 构建结果
            result = {
                "content": output_text,            # 文本内容
                "total_pages": total_pages,        # 总页数
                "total_lines": total_lines,        # 总行数
                "start_line": start_line_display,  # 起始行号
                "output_lines": truncation.output_lines  # 输出行数
            }
            
            if details:
                result["details"] = details
            
            return ToolResult.success(result)
            
        except Exception as e:
            return ToolResult.fail(f"Error reading PDF file: {str(e)}")
